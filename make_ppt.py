import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

prs = Presentation()
# 16:9 Widescreen slides
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

blank_layout = prs.slide_layouts[6]

def set_dark_bg(slide):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(11, 15, 25) # #0B0F19

# -------------------------------------------------------------
# SLIDE 1: Title Slide
# -------------------------------------------------------------
slide1 = prs.slides.add_slide(blank_layout)
set_dark_bg(slide1)

txBox = slide1.shapes.add_textbox(Inches(1), Inches(1.5), Inches(11.333), Inches(4.5))
tf = txBox.text_frame
tf.word_wrap = True

p1 = tf.paragraphs[0]
p1.text = "🌱 StorySprout"
p1.font.size = Pt(54)
p1.font.bold = True
p1.font.color.rgb = RGBColor(0, 242, 254) # Cyan gradient
p1.alignment = PP_ALIGN.CENTER

p2 = tf.add_paragraph()
p2.text = "Intergenerational Knowledge Transfer Through AI Storytelling — Turned Into a Book for Every Child"
p2.font.size = Pt(22)
p2.font.color.rgb = RGBColor(148, 163, 184) # Muted slate
p2.alignment = PP_ALIGN.CENTER
p2.space_before = Pt(20)

p3 = tf.add_paragraph()
p3.text = '"Sowing Yesterday\'s Wisdom, Growing Tomorrow\'s Minds."'
p3.font.size = Pt(20)
p3.font.italic = True
p3.font.color.rgb = RGBColor(245, 158, 11) # Amber Gold
p3.alignment = PP_ALIGN.CENTER
p3.space_before = Pt(30)

# -------------------------------------------------------------
# SLIDE 2: Architecture Diagram Slide
# -------------------------------------------------------------
slide2 = prs.slides.add_slide(blank_layout)
set_dark_bg(slide2)

# Title
txBox2 = slide2.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.333), Inches(0.8))
tf2 = txBox2.text_frame
p_title = tf2.paragraphs[0]
p_title.text = "Multi-Agent System Architecture & Model Flow"
p_title.font.size = Pt(28)
p_title.font.bold = True
p_title.font.color.rgb = RGBColor(0, 242, 254)

# Embed Diagram Image
img_path = r"c:\Users\jothi\Downloads\IBM_Projects\StorySprout\storysprout_architecture_diagram.png"
if os.path.exists(img_path):
    slide2.shapes.add_picture(img_path, Inches(0.5), Inches(1.2), width=Inches(12.333))

# -------------------------------------------------------------
# SLIDE 3: Verified Models Table Slide
# -------------------------------------------------------------
slide3 = prs.slides.add_slide(blank_layout)
set_dark_bg(slide3)

txBox3 = slide3.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.333), Inches(0.8))
tf3 = txBox3.text_frame
p_title3 = tf3.paragraphs[0]
p_title3.text = "Verified AI Models & Engine Specifications"
p_title3.font.size = Pt(28)
p_title3.font.bold = True
p_title3.font.color.rgb = RGBColor(0, 242, 254)

# Add Table
rows, cols = 8, 3
left, top, width, height = Inches(0.5), Inches(1.4), Inches(12.333), Inches(5.5)
table_shape = slide3.shapes.add_table(rows, cols, left, top, width, height)
table = table_shape.table

# Set Column Widths
table.columns[0].width = Inches(3.2)
table.columns[1].width = Inches(3.8)
table.columns[2].width = Inches(5.333)

headers = ["Component / Sub-Agent", "Verified Model / Engine", "System Function"]
data = [
    ["👑 Root Orchestrator", "Python ThreadPoolExecutor", "Coordinates sub-agent pipeline execution & thread concurrency"],
    ["🛡️ Safety Agent", "ibm/granite-4-h-small + Guardrails", "Sanitizes input prompts & page safety auditing"],
    ["✍️ Master Storyteller Agent", "ibm/granite-4-h-small", "Generates native narrative pages & title per age & language"],
    ["🔍 Fact Check Agent", "ibm/granite-4-h-small", "Validates historical dates, people & cultural traditions"],
    ["📖 Pedagogy Agent", "ibm/granite-4-h-small", "Extracts 4 age-tailored vocabulary words & meanings"],
    ["🧩 Quiz Agent", "ibm/granite-4-h-small", "Builds 3 comprehension multiple-choice questions & key"],
    ["🎨 Visual Director Agent", "ibm/granite-4-h-small", "Crafts scene illustration prompts & hero character sheet"],
]

for col_idx, h in enumerate(headers):
    cell = table.cell(0, col_idx)
    cell.fill.solid()
    cell.fill.fore_color.rgb = RGBColor(30, 41, 59)
    p = cell.text_frame.paragraphs[0]
    p.text = h
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = RGBColor(252, 211, 77)

for row_idx, row_data in enumerate(data, start=1):
    for col_idx, val in enumerate(row_data):
        cell = table.cell(row_idx, col_idx)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(15, 23, 42)
        p = cell.text_frame.paragraphs[0]
        p.text = val
        p.font.size = Pt(12)
        p.font.color.rgb = RGBColor(226, 232, 240)

# Save Presentation
output_ppt = r"c:\Users\jothi\Downloads\IBM_Projects\StorySprout\StorySprout_Architecture_Presentation.pptx"
prs.save(output_ppt)
print("PPT Presentation generated successfully at:", output_ppt)
